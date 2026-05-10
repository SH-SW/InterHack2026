"""
Build the "Heuristic Design — Cases We Considered" deck (English).

Walks through the design decisions: for each case, what we observed,
what we built, and why. Mythos brand: navy + steel-blue accents.
"""
from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]

# ---- Mythos brand palette (matches the other deck) ----
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
STEEL  = RGBColor(0x5A, 0x7A, 0x9A)
LIGHT  = RGBColor(0xC9, 0xD4, 0xE3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GREY   = RGBColor(0x6B, 0x6B, 0x6B)
SOFT   = RGBColor(0xF4, 0xF6, 0xFA)
ACCENT = RGBColor(0x2E, 0x6F, 0xA8)
WARN   = RGBColor(0xCC, 0x66, 0x33)

H_FONT = "Calibri"
B_FONT = "Calibri"


def fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, colour=INK,
             font=B_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = colour
    return box


def slide_title(slide, num, title):
    """Top number-tag + title."""
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.6), Inches(0.5),
                                    Inches(0.7), Inches(0.7))
    fill(badge, ACCENT)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(0.7), Inches(0.6),
             num, size=20, bold=True, colour=WHITE, font=H_FONT,
             align=PP_ALIGN.CENTER)
    # Title
    add_text(slide, Inches(1.5), Inches(0.55), Inches(11.3), Inches(0.7),
             title, size=28, bold=True, colour=NAVY, font=H_FONT)


def case_slide(prs, num, title, observed, decision, why):
    """
    Standard case slide layout:
      - Top: number badge + title
      - Two cards: "What we observed" (light) and "What we built" (navy)
      - Bottom strip: "Why this works"
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])

    slide_title(s, num, title)

    # LEFT card — What we observed
    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.6), Inches(1.7),
                              Inches(6.0), Inches(4.0))
    fill(left, SOFT); left.line.fill.background()
    add_text(s, Inches(0.85), Inches(1.85), Inches(5.5), Inches(0.4),
             "WHAT WE OBSERVED", size=11, bold=True, colour=WARN, font=H_FONT)
    add_text(s, Inches(0.85), Inches(2.30), Inches(5.5), Inches(3.2),
             observed, size=14, colour=INK, font=B_FONT)

    # Arrow between
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Inches(6.75), Inches(3.45),
                                Inches(0.5), Inches(0.4))
    fill(arrow, ACCENT)

    # RIGHT card — What we built
    right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(7.4), Inches(1.7),
                               Inches(5.3), Inches(4.0))
    fill(right, NAVY); right.line.fill.background()
    add_text(s, Inches(7.65), Inches(1.85), Inches(4.8), Inches(0.4),
             "WHAT WE BUILT", size=11, bold=True, colour=LIGHT, font=H_FONT)
    add_text(s, Inches(7.65), Inches(2.30), Inches(4.8), Inches(3.2),
             decision, size=14, colour=WHITE, font=B_FONT)

    # Bottom — Why
    add_text(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.4),
             "WHY", size=11, bold=True, colour=ACCENT, font=H_FONT)
    add_text(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(1.0),
             why, size=14, bold=True, colour=NAVY, font=B_FONT,
             align=PP_ALIGN.LEFT)


# ---- Build deck ----
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ====================================================================
# SLIDE 1 — Cover
# ====================================================================
s = prs.slides.add_slide(blank)
add_text(s, Inches(0.8), Inches(2.4), Inches(11), Inches(1.0),
         "Heuristic Design",
         size=64, bold=True, colour=NAVY, font=H_FONT)
add_text(s, Inches(0.8), Inches(3.4), Inches(11), Inches(1.0),
         "Cases We Considered",
         size=44, bold=False, colour=ACCENT, font=H_FONT)
add_text(s, Inches(0.8), Inches(4.7), Inches(11), Inches(0.5),
         "How and why each rule in the system was chosen — no black box, no overengineering.",
         size=18, colour=GREY, font=B_FONT)
# Bottom-right brand
add_text(s, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
         "Mythos Group  ·  Inibsa Alert Marker  ·  InterHack 2026",
         size=14, colour=GREY, font=B_FONT)
# Decorative accent dot
dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          Inches(11.3), Inches(2.5),
                          Inches(1.0), Inches(1.0))
fill(dot, ACCENT)


# ====================================================================
# SLIDE 2 — The principle
# ====================================================================
s = prs.slides.add_slide(blank)
slide_title(s, "00", "The principle")

add_text(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.6),
         "The brief is explicit:",
         size=18, colour=GREY, font=B_FONT)
add_text(s, Inches(0.6), Inches(2.25), Inches(12.1), Inches(0.7),
         "“The key is not technical sophistication, but analytical and commercial utility.”",
         size=22, bold=True, colour=NAVY, font=H_FONT)

# Three pillars
pillars = [
    ("Explainable",
     "Every alert is traceable to its inputs. Sales reps and judges can audit any decision."),
    ("Tunable",
     "Thresholds are heuristics by design — the learning loop adjusts them when real outcomes accumulate."),
    ("Anchored",
     "Each rule is grounded in either an observation in the data or an explicit requirement of the brief."),
]
y = Inches(3.8)
for i, (head, body) in enumerate(pillars):
    x = Inches(0.6 + i * 4.2)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              x, y, Inches(4.0), Inches(2.6))
    fill(card, SOFT); card.line.fill.background()
    add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(3.6), Inches(0.5),
             head, size=20, bold=True, colour=ACCENT, font=H_FONT)
    add_text(s, x + Inches(0.3), y + Inches(1.0), Inches(3.6), Inches(1.5),
             body, size=13, colour=INK, font=B_FONT)


# ====================================================================
# SLIDE 3 — Case 01: Two-track logic
# ====================================================================
case_slide(prs, "01", "Two-track logic — commodities vs. technical",
    observed=(
        "Two product blocks behave very differently in the data:\n\n"
        "• Commodities (anaesthesia, biosecurity) — recurring, predictable. Most clinics buy every month.\n\n"
        "• Technical products (biomaterials) — sporadic. A clinic may buy 4 times a year or once every 18 months."
    ),
    decision=(
        "Two parallel detection logics:\n\n"
        "• Commodities → share-of-potential rule against a declared annual potential.\n\n"
        "• Technical → individual-baseline rule comparing each clinic to its OWN past pattern."
    ),
    why=(
        "Applying one rule to both would either flag normal sporadic buyers as churning, "
        "or miss real loss in commodities. The brief insists on this distinction; the data confirms it."
    ),
)


# ====================================================================
# SLIDE 4 — Case 02: Seasonality
# ====================================================================
case_slide(prs, "02", "The August drop — seasonality is real",
    observed=(
        "Empirical seasonal factors per month per family:\n\n"
        "• August: 0.25–0.38× of the yearly average\n"
        "• November: 1.10–1.23× (peak)\n"
        "• Christmas week: also depressed\n\n"
        "Without seasonality awareness, ~1,500 false-positive churn alerts would fire every August."
    ),
    decision=(
        "Year-over-year baseline for technical products — compare to the SAME calendar window last year, "
        "not just the trailing 90 days.\n\n"
        "Plus an is_holiday_period() flag — drop signals are suppressed during August / Christmas / Reyes "
        "unless YoY also shows ≥50% decline."
    ),
    why=(
        "Comparing August to the previous spring is comparing apples to oranges. "
        "August is below-average for everyone — only same-vs-same comparisons isolate real churn."
    ),
)


# ====================================================================
# SLIDE 5 — Case 03: Lost-client filter
# ====================================================================
case_slide(prs, "03", "“Lost client” is more than 270 days silent — but not that simple",
    observed=(
        "Naive rule: client silent for >270 days = lost.\n\n"
        "Problem: lots of clients only ever made 1 or 2 purchases years ago. They were never really cyclic.\n\n"
        "Calling them “lost” triggers expensive recovery campaigns for clients who never engaged."
    ),
    decision=(
        "Cyclic-client filter on `lost` alerts. Three conditions, all required:\n\n"
        "• ≥3 distinct purchase days lifetime\n"
        "• ≥€200 lifetime volume\n"
        "• Cycle ≤365 days (or annualised volume ≥€500)\n\n"
        "Below those, no `lost` alert fires."
    ),
    why=(
        "A meaningful “lost” needs a meaningful “was active” first. "
        "The filter cuts the noise without losing real recovery candidates."
    ),
)


# ====================================================================
# SLIDE 6 — Case 04: Post-campaign grace
# ====================================================================
case_slide(prs, "04", "After a campaign, a quiet inbox doesn't mean churn",
    observed=(
        "When a clinic buys heavily during a promotional campaign, their warehouse is full.\n\n"
        "They will not buy again for several months — by design. That's the whole point of stocking up.\n\n"
        "Without context, the system would flag them as “gone silent” the next month."
    ),
    decision=(
        "Post-campaign grace period.\n\n"
        "If a client bought ≥2× their normal monthly average during a campaign window, "
        "suppress silent / churn-risk alerts for a period proportional to volume "
        "(volume / monthly_avg = grace months)."
    ),
    why=(
        "Otherwise the system pesters customers who already over-bought. "
        "Grace periods turn data we already have (Campañas table) into a real-world signal."
    ),
)


# ====================================================================
# SLIDE 7 — Case 05: Snooze
# ====================================================================
case_slide(prs, "05", "Don't fire the same alert every single day",
    observed=(
        "An alert is a state, not an event. The same client × family × type "
        "would fire the same alert every morning, indefinitely.\n\n"
        "Sales team would tune it out. Or worse, contact the same client repeatedly."
    ),
    decision=(
        "Outcome-driven snooze.\n\n"
        "When a sales rep records a terminal outcome (won, lost, or false_positive), "
        "that combination of `(client × family × tipo)` is suppressed for 30 days.\n\n"
        "Pending or no-contact outcomes do NOT snooze — those clients keep firing."
    ),
    why=(
        "The system should respect actions taken. Snoozing closes the loop between "
        "the activation layer and the feedback layer — the sales rep's work matters."
    ),
)


# ====================================================================
# SLIDE 8 — Case 06: Dynamic contact window
# ====================================================================
case_slide(prs, "06", "Contact window adapts to each client",
    observed=(
        "A fixed 7/30/90-day contact window is too rigid.\n\n"
        "A clinic with a 10-day cycle needs urgency now. A clinic with a 120-day cycle can wait.\n\n"
        "The right window depends on how soon the client is expected to buy."
    ),
    decision=(
        "Dynamic contact window per alert.\n\n"
        "If overdue → window collapses to default / 4 (act fast).\n\n"
        "If on-cycle → window = days until expected purchase.\n\n"
        "If no cycle data → fallback to default."
    ),
    why=(
        "An adaptive window respects the client's actual rhythm. "
        "Sales team prioritises by urgency that's real to the client, not arbitrary."
    ),
)


# ====================================================================
# SLIDE 9 — Case 07: Realistic capture impact
# ====================================================================
case_slide(prs, "07", "Don't claim €82K from a clinic that buys €5K/year",
    observed=(
        "Naive impact for a promiscuous client: full gap to declared potential.\n\n"
        "Example: declared potential €88K, current €5K → claim €82K of capturable demand.\n\n"
        "Reality: nobody captures 17× their current volume in one go. The number is indefensible."
    ),
    decision=(
        "Realistic capture impact = current volume × growth multiplier:\n\n"
        "• share <30% → impact = current × 0.5 (50% growth)\n"
        "• share ≥30% → impact = current × 0.3 (30% growth)\n\n"
        "Score remains driven by impact, just with believable numbers."
    ),
    why=(
        "Defensible numbers survive scrutiny. The brief asks for actionable alerts — "
        "and a sales rep can actually deliver 30-50% growth from a known active customer."
    ),
)


# ====================================================================
# SLIDE 10 — Case 08: Conversion probability
# ====================================================================
case_slide(prs, "08", "All four prioritisation signals from the brief in the score",
    observed=(
        "The brief lists four signals: economic impact, time urgency, client potential value, "
        "and conversion probability.\n\n"
        "Most teams score on `impact × urgency` only. Conversion probability silently disappears.\n\n"
        "The result: high-impact-but-unrecoverable alerts dominate, even though sales can't act on them."
    ),
    decision=(
        "Conversion probability per segment, used in scoring:\n\n"
        "  score = impact × urgency × conversion_probability\n\n"
        "Promiscuous: share / 0.30 (closer to loyal = more recoverable).\n"
        "Silent: decays with recency. Lost: 0.10. Spike: 0.50."
    ),
    why=(
        "The score now rewards recoverable opportunities, not just big numbers. "
        "All four brief signals influence the ranking — none silently dropped."
    ),
)


# ====================================================================
# SLIDE 11 — Summary
# ====================================================================
s = prs.slides.add_slide(blank)
slide_title(s, "✓", "Every rule has a reason")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.5),
         "Eight design cases. Each one anchored to either an observation in the data or a requirement in the brief.",
         size=15, colour=GREY, font=B_FONT)

# Two-column summary table
cases = [
    ("01  Two-track logic",        "Brief requirement"),
    ("02  Seasonality / August",   "Empirical (factor 0.25–0.38)"),
    ("03  Lost-client filter",     "Data quality / signal-to-noise"),
    ("04  Post-campaign grace",    "Brief interpretive caveat"),
    ("05  Outcome-driven snooze",  "Brief operational requirement"),
    ("06  Dynamic contact window", "Per-client rhythm in the data"),
    ("07  Realistic capture impact","Defensibility / business sense"),
    ("08  Conversion probability", "Brief: 4th prioritisation signal"),
]
y = Inches(2.5)
cw_l, cw_r = Inches(7.0), Inches(5.5)
for i, (label, anchor) in enumerate(cases):
    row_y = y + Inches(0.55 * i)
    # Left column — case name
    add_text(s, Inches(0.8), row_y, cw_l, Inches(0.5),
             label, size=15, bold=True, colour=NAVY, font=B_FONT)
    # Right column — anchored to
    add_text(s, Inches(7.6), row_y, cw_r, Inches(0.5),
             "→  " + anchor, size=14, colour=GREY, font=B_FONT)

# Closing line
add_text(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4),
         "No black box. Every threshold and multiplier is documented in the code.",
         size=14, bold=True, colour=ACCENT, font=B_FONT, align=PP_ALIGN.CENTER)


# ---- Save ----
out = ROOT / "pitch" / "Heuristic_Design_Cases.pptx"
prs.save(out)
print(f"✅ Wrote {out}  ({out.stat().st_size//1024} KB)  ·  {len(prs.slides)} slides")
