"""
Inibsa Alert Marker pitch deck — 8 slides following the official guide.

Section mapping:
  Slide 1  Cover (Mythos)
  Slide 2  Our approach           (En què consisteix el plantejament)
  Slide 3  Analytical logic       (Lògica analítica)
  Slide 4  Data, vars, assumptions (Dades, variables, supòsits)
  Slide 5  Output                 (Quin output genera)
  Slide 6  Technical info         (Info tècnica)
  Slide 7  Day-to-day operation   (Com funcionaria a la pràctica)
  Slide 8  Numbers + Thanks
"""
from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))

# ---- Mythos brand palette ----
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
STEEL  = RGBColor(0x5A, 0x7A, 0x9A)
LIGHT  = RGBColor(0xC9, 0xD4, 0xE3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1A, 0x1A, 0x1A)
GREY   = RGBColor(0x6B, 0x6B, 0x6B)
SOFT   = RGBColor(0xF4, 0xF6, 0xFA)
ACCENT = RGBColor(0x2E, 0x6F, 0xA8)
WARN   = RGBColor(0xCC, 0x66, 0x33)

H_FONT = "Calibri"
B_FONT = "Calibri"


def fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, colour=INK,
             font=B_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = colour
    return box


def add_bullets(slide, x, y, w, h, items, *, size=14, colour=INK, dot=ACCENT, spacing=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        r1 = p.add_run(); r1.text = "•  "; r1.font.size = Pt(size + 2)
        r1.font.bold = True; r1.font.color.rgb = dot; r1.font.name = B_FONT
        r2 = p.add_run(); r2.text = item; r2.font.size = Pt(size)
        r2.font.color.rgb = colour; r2.font.name = B_FONT
    return box


def slide_title(slide, num, title, subtitle=None):
    """Section number badge + title, optional subtitle line."""
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.6), Inches(0.5),
                                    Inches(0.7), Inches(0.7))
    fill(badge, ACCENT)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(0.7), Inches(0.6),
             num, size=20, bold=True, colour=WHITE, font=H_FONT,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.5), Inches(0.45), Inches(11.3), Inches(0.7),
             title, size=28, bold=True, colour=NAVY, font=H_FONT)
    if subtitle:
        add_text(slide, Inches(1.5), Inches(1.05), Inches(11.3), Inches(0.4),
                 subtitle, size=14, colour=GREY, font=B_FONT)


def fmt_eur(n):
    if n >= 1e6: return f"€{n/1e6:.1f}M"
    if n >= 1e3: return f"€{n/1e3:.0f}K"
    return f"€{n:,.0f}"


# ---- Live numbers ----
def _numbers():
    from smart_demand_signals import generate_alerts
    a = generate_alerts("2025-12-29")
    return {
        "n_alerts":     len(a),
        "n_high":       int((a["prioridad"] == "High").sum()),
        "total_impact": float(a["expected_impact_eur"].sum()),
    }


NUMS = _numbers()


# ---- Build deck ----
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ====================================================================
# SLIDE 1 — COVER (Mythos brand)
# ====================================================================
s = prs.slides.add_slide(blank)
cover_path = ROOT / "pitch" / "cover.png"
if cover_path.exists():
    s.shapes.add_picture(str(cover_path), 0, 0, width=Inches(13.333), height=Inches(7.5))
else:
    logo_x = Inches(2.0); logo_y = Inches(2.0); logo_size = Inches(2.6)
    left_v = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, logo_x, logo_y,
                                Inches(0.40), logo_size); fill(left_v, NAVY)
    right_v = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  logo_x + logo_size - Inches(0.40), logo_y,
                                  Inches(0.40), logo_size); fill(right_v, NAVY)
    top_l = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                logo_x, logo_y, Inches(0.95), Inches(0.40))
    fill(top_l, NAVY)
    top_r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                logo_x + logo_size - Inches(0.95), logo_y,
                                Inches(0.95), Inches(0.40))
    fill(top_r, NAVY)
    bars_x_start = logo_x + Inches(0.55)
    bars_y_base = logo_y + logo_size - Inches(0.20)
    bar_w = Inches(0.30); gap = Inches(0.05)
    heights = [0.45, 0.75, 1.05, 1.35]
    colours = [LIGHT, STEEL, STEEL, NAVY]
    for i, h in enumerate(heights):
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            bars_x_start + i * (bar_w + gap),
            bars_y_base - Inches(h), bar_w, Inches(h))
        fill(bar, colours[i])
    add_text(s, Inches(1.4), logo_y + logo_size + Inches(0.40),
             Inches(3.8), Inches(0.6),
             "MYTHOS", size=32, bold=False, colour=NAVY,
             font=H_FONT, align=PP_ALIGN.CENTER)
    underline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(2.6), logo_y + logo_size + Inches(1.05),
                                    Inches(1.4), Inches(0.04))
    fill(underline, NAVY)
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(6.2), Inches(1.7),
                              Inches(0.02), Inches(4.0))
    fill(div, GREY)
    add_text(s, Inches(6.7), Inches(2.6), Inches(6.4), Inches(1.0),
             "Inibsa Alert", size=54, bold=True, colour=INK,
             font=H_FONT, align=PP_ALIGN.LEFT)
    add_text(s, Inches(6.7), Inches(3.4), Inches(6.4), Inches(1.0),
             "Marker", size=54, bold=True, colour=INK,
             font=H_FONT, align=PP_ALIGN.LEFT)
    add_text(s, Inches(6.7), Inches(4.5), Inches(6.4), Inches(0.5),
             "Mythos Group", size=20, colour=GREY, font=B_FONT)
    add_text(s, Inches(6.7), Inches(4.95), Inches(6.4), Inches(0.5),
             "InterHack 2026", size=20, colour=GREY, font=B_FONT)


# ====================================================================
# SLIDE 2 — OUR APPROACH
# ====================================================================
s = prs.slides.add_slide(blank)
slide_title(s, "01", "Our approach",
            "Central idea, why this method, and the advantage over alternatives")

# Central idea card (full width)
idea = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(1.65),
                          Inches(12.13), Inches(1.6))
fill(idea, NAVY); idea.line.fill.background()
add_text(s, Inches(0.9), Inches(1.85), Inches(11.6), Inches(0.4),
         "CENTRAL IDEA", size=11, bold=True, colour=LIGHT, font=H_FONT)
add_text(s, Inches(0.9), Inches(2.20), Inches(11.6), Inches(1.0),
         "Turn 5 years of Inibsa's sales history into a daily ranked call list — "
         "every clinic that needs attention today, with a clear reason and a recommended action.",
         size=18, bold=True, colour=WHITE, font=B_FONT)

# Two cards: Why this method  +  Advantage
y = Inches(3.55)
ch = Inches(3.4)
left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), y, Inches(5.95), ch)
fill(left, SOFT); left.line.fill.background()
add_text(s, Inches(0.9), y + Inches(0.2), Inches(5.5), Inches(0.4),
         "WHY HEURISTICS, NOT ML", size=11, bold=True, colour=ACCENT, font=H_FONT)
add_bullets(s, Inches(0.9), y + Inches(0.65), Inches(5.5), ch - Inches(0.8), [
    "Brief explicitly says “the key is analytical and commercial utility, not technical sophistication.”",
    "ML needs labelled outcomes — we don't have them yet, only mocked.",
    "Rules are tunable later via the feedback loop, when real outcomes accumulate.",
], size=12)

right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(6.78), y, Inches(5.95), ch)
fill(right, SOFT); right.line.fill.background()
add_text(s, Inches(7.08), y + Inches(0.2), Inches(5.5), Inches(0.4),
         "OUR ADVANTAGE", size=11, bold=True, colour=ACCENT, font=H_FONT)
add_bullets(s, Inches(7.08), y + Inches(0.65), Inches(5.5), ch - Inches(0.8), [
    "Explainable: every alert traces to its raw inputs. Sales reps and judges can audit.",
    "Defensible: every threshold is anchored to either an observation or a brief requirement.",
    "Fast: 5 seconds for 6,000 clients. Daily idempotent. Linear scaling.",
], size=12)


# ====================================================================
# SLIDE 3 — ANALYTICAL LOGIC
# ====================================================================
s = prs.slides.add_slide(blank)
slide_title(s, "02", "Analytical logic",
            "Two parallel logics · signals per case · noise filtering · priority scoring")

# Top: Two-track block
top_y = Inches(1.65)
left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), top_y, Inches(5.95), Inches(2.6))
fill(left, SOFT); left.line.fill.background()
add_text(s, Inches(0.9), top_y + Inches(0.15), Inches(5.5), Inches(0.4),
         "COMMODITIES — anaesthesia · biosecurity", size=11, bold=True,
         colour=ACCENT, font=H_FONT)
add_text(s, Inches(0.9), top_y + Inches(0.55), Inches(5.5), Inches(0.5),
         "Recurring purchases · share-of-potential rule",
         size=14, bold=True, colour=NAVY, font=H_FONT)
add_bullets(s, Inches(0.9), top_y + Inches(1.05), Inches(5.5), Inches(1.5), [
    "Loyal ≥30%, Promiscuous 5-30%, Marginal <5%",
    "Lost: >270d silent + cyclic-client filter",
    "Trend: 6m current vs prior period",
], size=11, spacing=4)

right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(6.78), top_y, Inches(5.95), Inches(2.6))
fill(right, NAVY); right.line.fill.background()
add_text(s, Inches(7.08), top_y + Inches(0.15), Inches(5.5), Inches(0.4),
         "TECHNICAL — biomaterials", size=11, bold=True, colour=LIGHT, font=H_FONT)
add_text(s, Inches(7.08), top_y + Inches(0.55), Inches(5.5), Inches(0.5),
         "Sporadic purchases · individual baseline + YoY",
         size=14, bold=True, colour=WHITE, font=H_FONT)
add_bullets(s, Inches(7.08), top_y + Inches(1.05), Inches(5.5), Inches(1.5), [
    "Compare to client's OWN past pattern",
    "4 signals: absent · drop_freq · drop_volume · anomaly",
    "Year-over-year baseline neutralises August / Christmas",
], size=11, colour=WHITE, dot=LIGHT, spacing=4)

# Bottom: Noise filtering + Priority scoring
bot_y = Inches(4.45)
nl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6), bot_y, Inches(7.4), Inches(2.55))
fill(nl, SOFT); nl.line.fill.background()
add_text(s, Inches(0.9), bot_y + Inches(0.2), Inches(6.9), Inches(0.4),
         "NOISE → SIGNAL", size=11, bold=True, colour=WARN, font=H_FONT)
add_bullets(s, Inches(0.9), bot_y + Inches(0.65), Inches(6.9), Inches(1.8), [
    "Holiday flag suppresses August/Christmas drops unless YoY confirms",
    "Post-campaign grace silences clients who just stocked up heavily",
    "Outcome-driven snooze: 30d after won/lost/false_positive on same (cliente×fam×tipo)",
    "Cyclic-client filter: lost only fires if client was ever truly cyclic",
], size=11, spacing=3)

pr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(8.23), bot_y, Inches(4.5), Inches(2.55))
fill(pr, SOFT); pr.line.fill.background()
add_text(s, Inches(8.53), bot_y + Inches(0.2), Inches(4.0), Inches(0.4),
         "PRIORITY = SCORE", size=11, bold=True, colour=ACCENT, font=H_FONT)
add_text(s, Inches(8.53), bot_y + Inches(0.65), Inches(4.0), Inches(1.0),
         "score = impact × urgency × conv_prob",
         size=13, bold=True, colour=NAVY, font="Consolas")
add_text(s, Inches(8.53), bot_y + Inches(1.30), Inches(4.0), Inches(1.2),
         "All four prioritisation signals from the brief.\n\n"
         "Higher score = higher in the daily call list.",
         size=11, colour=INK, font=B_FONT)


# ====================================================================
# SLIDE 4 — DATA, VARIABLES, ASSUMPTIONS
# ====================================================================
s = prs.slides.add_slide(blank)
slide_title(s, "03", "Data, variables, assumptions",
            "Datasets used · derived variables · what we assume · what we cannot observe")

# Three columns
col_y = Inches(1.65)
col_h = Inches(5.6)
cw = Inches(4.0)
gap = Inches(0.13)

# Col 1 — Datasets + Direct vars
c1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6), col_y, cw, col_h)
fill(c1, SOFT); c1.line.fill.background()
add_text(s, Inches(0.85), col_y + Inches(0.2), cw - Inches(0.5), Inches(0.4),
         "DATASETS USED", size=11, bold=True, colour=ACCENT, font=H_FONT)
add_bullets(s, Inches(0.85), col_y + Inches(0.6), cw - Inches(0.5), Inches(2.3), [
    "Ventas — 162K rows, 5 yrs, daily granularity",
    "Clientes — 11K, postal code, province",
    "Productos — 25 SKUs, 4 sub-families",
    "Potencial — annual estimate per client × familia",
    "Campañas — 10 promotional windows",
], size=11, spacing=2)
add_text(s, Inches(0.85), col_y + Inches(3.25), cw - Inches(0.5), Inches(0.4),
         "DIRECT VARIABLES", size=11, bold=True, colour=ACCENT, font=H_FONT)
add_bullets(s, Inches(0.85), col_y + Inches(3.65), cw - Inches(0.5), Inches(1.9), [
    "fecha · unidades · valores_h",
    "id_cliente · id_producto",
    "potencial_h · provincia",
    "tipo_transaccion (cleaned)",
], size=11, spacing=2)

# Col 2 — Derived variables
c2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6) + cw + gap, col_y, cw, col_h)
fill(c2, NAVY); c2.line.fill.background()
add_text(s, Inches(0.85) + cw + gap, col_y + Inches(0.2),
         cw - Inches(0.5), Inches(0.4),
         "DERIVED VARIABLES", size=11, bold=True, colour=LIGHT, font=H_FONT)
add_bullets(s, Inches(0.85) + cw + gap, col_y + Inches(0.6),
            cw - Inches(0.5), col_h - Inches(0.8), [
    "share_of_potential (trailing 12m)",
    "mean_interpurchase_days",
    "cyclicity_score · cycle_progress",
    "recency_days · last_purchase",
    "freq_drop_ratio · vol_drop_ratio",
    "expected_freq_recent (YoY)",
    "expected_vol_recent (YoY)",
    "loyalty_tier · trend",
    "conversion_probability",
    "score = impact × urgency × conv_prob",
], size=11, colour=WHITE, dot=LIGHT, spacing=2)

# Col 3 — Assumptions + Inferred
c3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6) + 2 * (cw + gap), col_y, cw, col_h)
fill(c3, SOFT); c3.line.fill.background()
add_text(s, Inches(0.85) + 2 * (cw + gap), col_y + Inches(0.2),
         cw - Inches(0.5), Inches(0.4),
         "KEY ASSUMPTIONS", size=11, bold=True, colour=WARN, font=H_FONT)
add_bullets(s, Inches(0.85) + 2 * (cw + gap), col_y + Inches(0.6),
            cw - Inches(0.5), Inches(2.5), [
    "Cyclicity inferred from purchase intervals",
    "potencial_h is Inibsa's annual estimate (≥0)",
    "Returns/zero-cost units flagged, not dropped",
    "Unregistered clients flagged separately",
], size=11, spacing=2)
add_text(s, Inches(0.85) + 2 * (cw + gap), col_y + Inches(3.05),
         cw - Inches(0.5), Inches(0.4),
         "WHAT WE CAN'T OBSERVE", size=11, bold=True, colour=WARN, font=H_FONT)
add_text(s, Inches(0.85) + 2 * (cw + gap), col_y + Inches(3.45),
         cw - Inches(0.5), col_h - Inches(3.7),
         "Competitor purchases are invisible. We INFER un-captured demand "
         "from the gap between observed sales and Inibsa's potential estimate. "
         "share_of_potential <30% with active recent buying = promiscuous (split with competitors).",
         size=11, colour=INK, font=B_FONT)


# ====================================================================
# SLIDE 5 — OUTPUT
# ====================================================================
s = prs.slides.add_slide(blank)
slide_title(s, "04", "Output",
            "Alert structure · what's inside · how it's surfaced to the sales team")

# Left: Example alert (terminal-style)
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.6), Inches(1.65),
                           Inches(7.4), Inches(5.4))
fill(panel, INK); panel.line.fill.background()
add_text(s, Inches(0.85), Inches(1.85), Inches(7.0), Inches(0.4),
         "ONE ALERT — JSON-LIKE", size=11, bold=True, colour=WARN, font="Consolas")

fields = [
    ("alert_id",            '"ALT-20251229-000001"'),
    ("id_cliente",          '"40439"'),
    ("provincia",           '"Zaragoza"'),
    ("familia",             '"Categoria C2"  (Bioseguridad)'),
    ("tipo_alerta",         '"capture_window"'),
    ("category",            '"sales_opportunity"'),
    ("prioridad",           '"High"'),
    ("score",               '22,613'),
    ("expected_impact_eur", '€2,496  (current × 0.5)'),
    ("urgency_factor",      '0.70'),
    ("conversion_probability", '0.83  (near-loyal)'),
    ("canal",               '"delegado"'),
    ("contact_window_days", '6  (dynamic)'),
    ("motivo",              '"Promiscuous client: 27% share..."'),
    ("trace_features",      '{ share: 0.27, current: 4,993, ... }'),
]
y = Inches(2.30)
for k, v in fields:
    add_text(s, Inches(0.95), y, Inches(2.7), Inches(0.30),
             k, size=11, bold=True, colour=WARN, font="Consolas")
    add_text(s, Inches(3.65), y, Inches(4.2), Inches(0.30),
             v, size=11, colour=WHITE, font="Consolas")
    y += Inches(0.31)

# Right: How it's surfaced
add_text(s, Inches(8.25), Inches(1.85), Inches(4.5), Inches(0.4),
         "HOW IT'S SURFACED", size=11, bold=True, colour=ACCENT, font=H_FONT)

# Three surfaces with mini-cards
surfaces = [
    ("Streamlit dashboard",
     "3 tabs · Client view (4 categories) · Monitoring · Learning loop",
     ACCENT),
    ("CSV / JSON export",
     "Filtered call list · HubSpot Tasks-shaped JSON",
     STEEL),
    ("Daily call list",
     "🔴 Lost · 🟠 Loss risk · 🟡 Sales opportunity · 🟢 Healthy",
     NAVY),
]
sy = Inches(2.3)
for title, body, clr in surfaces:
    bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(8.25), sy, Inches(4.5), Inches(1.45))
    fill(bx, SOFT); bx.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(8.25), sy,
                              Inches(0.10), Inches(1.45))
    fill(bar, clr)
    add_text(s, Inches(8.5), sy + Inches(0.20), Inches(4.2), Inches(0.4),
             title, size=14, bold=True, colour=NAVY, font=H_FONT)
    add_text(s, Inches(8.5), sy + Inches(0.65), Inches(4.2), Inches(0.7),
             body, size=11, colour=INK, font=B_FONT)
    sy += Inches(1.6)

add_text(s, Inches(8.25), Inches(7.0), Inches(4.5), Inches(0.3),
         "Every metric in trace_features → no black box.",
         size=10, bold=True, colour=ACCENT, font=B_FONT)


# ====================================================================
# SLIDE 6 — TECHNICAL INFO
# ====================================================================
s = prs.slides.add_slide(blank)
slide_title(s, "05", "Technical info",
            "How we built it · language · environment · tooling")

# Three pillars
items = [
    ("HOW WE BUILT IT",
     [
        "AI-assisted (Claude Code) + manual review on every change",
        "13 PRs of clean git history",
        "21 smoke tests passing",
        "Cleaned dataset (5 sheets, 162K rows of Ventas)",
     ],
     ACCENT),
    ("LANGUAGE & STACK",
     [
        "Python 3.10+",
        "pandas, numpy — analytics",
        "streamlit — interactive dashboard",
        "openpyxl — clean Excel I/O",
        "python-pptx — programmatic deck",
     ],
     STEEL),
    ("ENVIRONMENT",
     [
        "macOS · Linux · Windows",
        "Single command:  streamlit run src/dashboard.py",
        "No internet required at runtime",
        "Data + code self-contained in repo",
        "≤5 sec daily compute for 6,000 clients",
     ],
     NAVY),
]
y = Inches(1.65)
ch = Inches(5.6)
cw = Inches(4.0)
gap = Inches(0.13)
for i, (title, rows, clr) in enumerate(items):
    x = Inches(0.6) + i * (cw + gap)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    fill(card, SOFT if i < 2 else NAVY); card.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.10), ch)
    fill(bar, clr)
    add_text(s, x + Inches(0.3), y + Inches(0.2), cw - Inches(0.5), Inches(0.4),
             title, size=11, bold=True,
             colour=clr if i < 2 else LIGHT, font=H_FONT)
    add_bullets(s, x + Inches(0.3), y + Inches(0.7),
                cw - Inches(0.5), ch - Inches(0.9),
                rows, size=12,
                colour=INK if i < 2 else WHITE,
                dot=ACCENT if i < 2 else LIGHT)


# ====================================================================
# SLIDE 7 — DAY-TO-DAY OPERATION
# ====================================================================
s = prs.slides.add_slide(blank)
slide_title(s, "06", "Day-to-day operation",
            "Daily workflow · minimum architecture · live demo")

# Left column: workflow
add_text(s, Inches(0.6), Inches(1.7), Inches(7.0), Inches(0.4),
         "DAILY WORKFLOW", size=11, bold=True, colour=ACCENT, font=H_FONT)
steps = [
    ("08:00", "Engine runs", "generate_alerts(today) — ~5 sec"),
    ("08:05", "Sales team opens dashboard", "Categories at a glance"),
    ("08:10", "Sales rep clicks a client", "Cycle status + reason + action"),
    ("Day", "Acts on alerts", "Phone / visit / email"),
    ("End", "Records outcome", "won / lost / pending / no_contact"),
]
sy = Inches(2.15)
for time, action, detail in steps:
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.6), sy, Inches(0.85), Inches(0.5))
    fill(pill, NAVY)
    add_text(s, Inches(0.6), sy + Inches(0.05), Inches(0.85), Inches(0.4),
             time, size=11, bold=True, colour=WHITE, font=H_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.6), sy + Inches(0.04), Inches(5.4), Inches(0.3),
             action, size=13, bold=True, colour=NAVY, font=H_FONT)
    add_text(s, Inches(1.6), sy + Inches(0.32), Inches(5.4), Inches(0.3),
             detail, size=11, colour=GREY, font=B_FONT)
    sy += Inches(0.62)

# Right column: minimum architecture
add_text(s, Inches(8.0), Inches(1.7), Inches(5.0), Inches(0.4),
         "MINIMUM ARCHITECTURE", size=11, bold=True, colour=ACCENT, font=H_FONT)

layers = [
    ("Feedback layer", "alert_outcomes.csv → metrics → tuning",   ACCENT),
    ("Activation layer", "build_alerts → ranked, traceable",      NAVY),
    ("Analytical layer", "commodity_segments / technical_patterns", STEEL),
    ("Data layer",     "cleaned CSVs (5 sheets)",                 GREY),
]
ly = Inches(2.15)
for title, sub, clr in layers:
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(8.0), ly, Inches(5.0), Inches(0.85))
    fill(box, clr); box.line.fill.background()
    add_text(s, Inches(8.2), ly + Inches(0.10), Inches(4.6), Inches(0.4),
             title, size=14, bold=True, colour=WHITE, font=H_FONT)
    add_text(s, Inches(8.2), ly + Inches(0.48), Inches(4.6), Inches(0.4),
             sub, size=10, colour=WHITE, font=B_FONT)
    ly += Inches(0.95)
    if title != "Data layer":
        arr = s.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                  Inches(10.3), ly - Inches(0.1),
                                  Inches(0.4), Inches(0.13))
        fill(arr, ACCENT)

# Demo callout
add_text(s, Inches(0.6), Inches(6.55), Inches(12.13), Inches(0.4),
         "🎬 Live demo: 3 minutes · pitch/demo_walkthrough_3min.md",
         size=14, bold=True, colour=NAVY, font=B_FONT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(6.95), Inches(12.13), Inches(0.35),
         "AI-narrated demo video also available · pitch/video/Inibsa_Alert_Marker_demo.mp4",
         size=11, colour=GREY, font=B_FONT, align=PP_ALIGN.CENTER)


# ====================================================================
# SLIDE 8 — NUMBERS + THANKS
# ====================================================================
s = prs.slides.add_slide(blank)
add_text(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.7),
         "By the numbers · 8 of 8 deliverables shipped",
         size=30, bold=True, colour=NAVY, font=H_FONT)

metrics = [
    (f"{NUMS['n_alerts']:,}", "alerts generated today"),
    (f"{NUMS['n_high']:,}",   "High priority"),
    (fmt_eur(NUMS['total_impact']), "expected commercial impact"),
]
y_metrics = Inches(1.6)
for i, (val, lbl) in enumerate(metrics):
    x = Inches(0.7 + i * 4.2)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              x, y_metrics, Inches(4.0), Inches(1.7))
    fill(card, SOFT); card.line.fill.background()
    add_text(s, x, y_metrics + Inches(0.25), Inches(4.0), Inches(0.9),
             val, size=44, bold=True, colour=NAVY, font=H_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y_metrics + Inches(1.20), Inches(4.0), Inches(0.4),
             lbl, size=13, colour=GREY, font=B_FONT, align=PP_ALIGN.CENTER)

deliverables = [
    "Purchase prediction (commodities)",
    "Early churn-risk (technical)",
    "Capture-window identification",
    "Interpretable, actionable alerts",
    "Operational prioritisation",
    "Daily-operation proposal",
    "Standalone-to-CRM evolution path",
    "Learning capability (feedback loop)",
]
y_sc = Inches(3.6)
for i, txt in enumerate(deliverables):
    c = i % 2; r = i // 2
    x = Inches(0.7 + c * 6.15)
    y = y_sc + r * Inches(0.65)
    tick = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.32), Inches(0.32))
    fill(tick, ACCENT)
    add_text(s, x, y, Inches(0.32), Inches(0.32),
             "✓", size=14, bold=True, colour=WHITE, font=H_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.50), y + Inches(0.04),
             Inches(5.5), Inches(0.4),
             txt, size=13, colour=INK, font=B_FONT)

ty = Inches(6.6)
add_text(s, Inches(0), ty, Inches(13.333), Inches(0.5),
         "Thank you · Questions?",
         size=18, bold=True, colour=NAVY, font=H_FONT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), ty + Inches(0.45), Inches(13.333), Inches(0.4),
         "Mythos Group · github.com/SH-SW/InterHack2026",
         size=11, colour=GREY, font=B_FONT, align=PP_ALIGN.CENTER)


# ---- Save ----
out = ROOT / "pitch" / "Inibsa_Alert_Marker.pptx"
prs.save(out)
print(f"✅ Wrote {out}  ({out.stat().st_size//1024} KB)  ·  {len(prs.slides)} slides")
print(f"   Live numbers: {NUMS['n_alerts']:,} alerts · {NUMS['n_high']:,} High · {fmt_eur(NUMS['total_impact'])}")
